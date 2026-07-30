from __future__ import annotations

import json
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE


BUNDLE_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(BUNDLE_DIR))

import builder  # noqa: E402
from validate_deck import inspect_deck  # noqa: E402


class LayoutSafetyTests(unittest.TestCase):
    def _slide(self):
        prs = Presentation()
        return prs.slides.add_slide(prs.slide_layouts[6])

    def test_textbox_rejects_content_that_cannot_fit(self):
        with self.assertRaisesRegex(ValueError, "does not fit"):
            builder._textbox(
                self._slide(),
                "Слишком длинный текст " * 40,
                0,
                0,
                2,
                0.45,
                size=20,
                color="#111111",
                font="Arial",
                label="regression text",
            )

    def test_textbox_enables_renderer_level_autofit(self):
        shape = builder._textbox(
            self._slide(),
            "Короткий текст",
            0,
            0,
            3,
            0.6,
            size=18,
            color="#111111",
            font="Arial",
        )
        self.assertEqual(shape.text_frame.auto_size, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)

    def test_bullet_list_rejects_excessive_density(self):
        with self.assertRaisesRegex(ValueError, "bullet list does not fit"):
            builder._add_paragraphs(
                self._slide(),
                ["Очень длинный пункт " * 25] * 5,
                0,
                0,
                3,
                1.2,
                size=18,
                color="#111111",
                accent="#3366FF",
                font="Arial",
            )

    def test_build_rejects_overloaded_slide_before_delivery(self):
        spec = {
            "filename": "overloaded.pptx",
            "theme": "paper",
            "slides": [
                {
                    "type": "title",
                    "title": "Очень длинный заголовок " * 24,
                },
                {
                    "type": "closing",
                    "title": "Следующий шаг",
                },
            ],
        }
        with tempfile.TemporaryDirectory(prefix="presentation-overflow-test-") as tmp:
            with self.assertRaisesRegex(ValueError, "does not fit"):
                builder.build_from_spec(spec, tmp)

    def test_showcase_build_passes_structural_qa(self):
        spec = json.loads(
            (BUNDLE_DIR / "examples" / "example_spec.json").read_text(
                encoding="utf-8"
            )
        )
        with tempfile.TemporaryDirectory(prefix="presentation-layout-test-") as tmp:
            path = Path(builder.build_from_spec(spec, tmp))
            report = inspect_deck(path)
        self.assertEqual(report["issues"], [])


if __name__ == "__main__":
    unittest.main()
