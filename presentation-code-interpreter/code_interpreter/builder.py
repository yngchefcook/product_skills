#!/usr/bin/env python3
"""Portable, data-driven PPTX builder for Code Interpreter.

Public API:
    ``lint_spec(spec)`` validates content and layout constraints without writing.
    ``validate_spec(spec)`` returns blocking errors for compatibility.
    ``build_from_spec(spec, out_dir)`` writes a presentation and returns its path.

Image fields accept a catalog ID, a local path, or an HTTP(S) ``image_url``.
Remote images are downloaded into a temporary directory, validated with Pillow,
embedded in the PPTX, and removed after the file is saved. See
``SPEC_REFERENCE.md`` for the complete input schema.

CLI example:
    python3 builder.py spec.json --output-dir ./output
"""

from __future__ import annotations

import copy
import math
import tempfile
import textwrap
from pathlib import Path
from typing import Any, Iterable, Sequence

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_DATA_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
)
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

try:  # Package import after uploading the enclosing folder.
    from .helpers import (
        THEMES_PATH,
        assert_bundle_complete,
        download_remote_image,
        is_http_url,
        load_json,
        sanitize_filename,
    )
    from .image_filenames import resolve_image as resolve_catalog_image
except ImportError:  # Direct ``python builder.py`` or sys.path-based import.
    from helpers import (
        THEMES_PATH,
        assert_bundle_complete,
        download_remote_image,
        is_http_url,
        load_json,
        sanitize_filename,
    )
    from image_filenames import resolve_image as resolve_catalog_image

__all__ = ["build_from_spec", "lint_spec", "validate_spec"]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
W_IN = 13.333
H_IN = 7.5
MARGIN = 0.78
BLANK_LAYOUT = 6

SLIDE_TYPES = frozenset(
    {
        "title",
        "section",
        "statement",
        "bullets",
        "image",
        "metrics",
        "chart",
        "comparison",
        "timeline",
        "process",
        "quote",
        "table",
        "summary",
        "closing",
    }
)
IMAGE_SLIDE_TYPES = frozenset({"title", "section", "image", "quote", "closing"})

REQUIRED_FIELDS: dict[str, tuple[str, ...]] = {
    "title": ("title",),
    "section": ("title",),
    "statement": ("title",),
    "bullets": ("title", "bullets"),
    "image": ("title",),
    "metrics": ("title", "metrics"),
    "chart": ("title", "chart_type", "categories", "series"),
    "comparison": ("title", "left", "right"),
    "timeline": ("title", "steps"),
    "process": ("title", "steps"),
    "quote": ("quote",),
    "table": ("title", "columns", "rows"),
    "summary": ("title", "takeaway", "actions"),
    "closing": ("title",),
}

CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}


def _load_json(path: Path) -> dict:
    return load_json(path)


def _load_theme(theme_value: str | dict | None) -> dict:
    """Resolve a named theme or merge ``overrides`` on top of a base theme."""
    data = _load_json(THEMES_PATH)
    themes = data["themes"]
    if theme_value is None:
        return dict(themes[data["default"]])
    if isinstance(theme_value, str):
        if theme_value not in themes:
            raise ValueError(
                f"Unknown theme {theme_value!r}. Available: {sorted(themes)}"
            )
        return dict(themes[theme_value])
    if isinstance(theme_value, dict):
        base_name = theme_value.get("base", data["default"])
        if base_name not in themes:
            raise ValueError(
                f"Unknown base theme {base_name!r}. Available: {sorted(themes)}"
            )
        result = dict(themes[base_name])
        result.update(theme_value.get("overrides", {}))
        return result
    raise TypeError("theme must be a theme name or an object with base/overrides")


def _rgb(value: str | Sequence[int]) -> RGBColor:
    if isinstance(value, str):
        raw = value.strip().lstrip("#")
        if len(raw) != 6:
            raise ValueError(f"Invalid color: {value!r}")
        return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
    if len(value) != 3:
        raise ValueError(f"Invalid RGB triplet: {value!r}")
    return RGBColor(int(value[0]), int(value[1]), int(value[2]))


def _set_opacity(shape, opacity: float) -> None:
    """Set solid-fill opacity using DrawingML (0.0 transparent, 1.0 opaque)."""
    opacity = max(0.0, min(1.0, float(opacity)))
    solid_fill = shape.fill._xPr.solidFill
    if solid_fill is None:
        return
    color = solid_fill.getchildren()[0] if len(solid_fill) else None
    if color is None:
        return
    for old in list(color):
        if old.tag.endswith("}alpha"):
            color.remove(old)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(int(opacity * 100000)))
    color.append(alpha)


def _rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    color: str | Sequence[int],
    *,
    opacity: float = 1.0,
    radius: bool = False,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    _set_opacity(shape, opacity)
    shape.line.fill.background()
    return shape


def _line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str | Sequence[int],
    *,
    width: float = 1.0,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width)
    return line


def _set_cell_margins(cell, left=0.08, right=0.08, top=0.05, bottom=0.05):
    cell.margin_left = Inches(left)
    cell.margin_right = Inches(right)
    cell.margin_top = Inches(top)
    cell.margin_bottom = Inches(bottom)


def _set_cell_borders(cell, color: str | Sequence[int], width_pt: float = 0.55):
    """Write table-cell borders directly in DrawingML for renderer parity."""
    if isinstance(color, str):
        hex_color = color.strip().lstrip("#").upper()
    else:
        hex_color = "".join(f"{int(channel):02X}" for channel in color)
    tc_pr = cell._tc.get_or_add_tcPr()
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        node = tc_pr.find(qn(edge))
        if node is None:
            node = OxmlElement(edge)
            tc_pr.append(node)
        node.set("w", str(int(width_pt * 12700)))
        node.set("cap", "flat")
        for child in list(node):
            node.remove(child)
        solid = OxmlElement("a:solidFill")
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", hex_color)
        solid.append(srgb)
        dash = OxmlElement("a:prstDash")
        dash.set("val", "solid")
        node.append(solid)
        node.append(dash)


def _wrapped_line_count(text: str, width_in: float, font_pt: float) -> int:
    """Conservatively estimate wrapped lines for a fixed-width text box."""
    # Average Latin/Cyrillic body glyphs occupy roughly 0.52 em. This matches
    # the structural validator while leaving renderer-level variance to
    # PowerPoint's text-to-fit safety setting.
    chars_per_line = max(1, int(width_in * 72 / max(font_pt * 0.52, 1)))
    lines = 0
    for raw_line in str(text).splitlines() or [""]:
        if not raw_line.strip():
            lines += 1
            continue
        lines += max(
            1,
            len(
                textwrap.wrap(
                    raw_line,
                    width=chars_per_line,
                    break_long_words=True,
                    break_on_hyphens=True,
                    replace_whitespace=False,
                )
            ),
        )
    return lines


def _assert_text_fits(
    text: str,
    w: float,
    h: float,
    *,
    size: float,
    margin: float = 0.0,
    line_spacing: float | None = None,
    label: str = "text box",
) -> None:
    """Raise before rendering when text cannot fit its assigned geometry."""
    usable_w = max(0.1, float(w) - 2 * float(margin))
    usable_h = max(0.05, float(h) - 2 * float(margin) - 0.03)
    required_lines = _wrapped_line_count(str(text), usable_w, float(size))
    spacing_factor = (
        max(1.0, float(line_spacing)) if line_spacing is not None else 1.16
    )
    line_height_in = float(size) / 72 * spacing_factor
    available_lines = max(1, math.floor(usable_h / max(line_height_in, 0.01)))
    if required_lines > available_lines:
        excerpt = " ".join(str(text).split())
        if len(excerpt) > 72:
            excerpt = f"{excerpt[:69]}..."
        raise ValueError(
            f"{label} does not fit: {required_lines} estimated lines, "
            f"{available_lines} available at {size:g} pt in {w:g}×{h:g} in. "
            f"Shorten the content or choose another slide type: {excerpt!r}"
        )


def _textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str | Sequence[int],
    font: str,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.0,
    line_spacing: float | None = None,
    label: str = "text box",
):
    _assert_text_fits(
        str(text),
        w,
        h,
        size=size,
        margin=margin,
        line_spacing=line_spacing,
        label=label,
    )
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    # Keep explicit font sizes as the design contract, while allowing Office to
    # make only renderer-level safety adjustments instead of overflowing into
    # the next object.
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return shape


def _add_paragraphs(
    slide,
    items: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str | Sequence[int],
    accent: str | Sequence[int],
    font: str,
    numbered: bool = False,
    max_items: int = 5,
    spacing: float = 10,
):
    rendered_items = [str(item) for item in list(items)[:max_items]]
    required_height = 0.0
    for item in rendered_items:
        lines = _wrapped_line_count(
            f"{'00  ' if numbered else '●  '}{item}",
            float(w),
            float(size),
        )
        required_height += lines * (float(size) / 72 * 1.16)
    if len(rendered_items) > 1:
        required_height += (len(rendered_items) - 1) * (float(spacing) / 72)
    if required_height > float(h) - 0.03:
        raise ValueError(
            f"bullet list does not fit: needs approximately "
            f"{required_height:.2f} in, {h:.2f} in available. "
            "Shorten bullets, reduce their count, or choose another slide type."
        )
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for index, item in enumerate(rendered_items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        p.line_spacing = 1.08
        mark = p.add_run()
        mark.text = f"{index + 1:02d}  " if numbered else "●  "
        mark.font.name = font
        mark.font.size = Pt(size - 1 if numbered else max(10, size - 4))
        mark.font.bold = numbered
        mark.font.color.rgb = _rgb(accent)
        body = p.add_run()
        body.text = str(item)
        body.font.name = font
        body.font.size = Pt(size)
        body.font.color.rgb = _rgb(color)
    return shape


def _image_ref(slide_spec: dict) -> str | Path | None:
    """Return the single image reference accepted by a slide specification."""
    return slide_spec.get("image") or slide_spec.get("image_url")


def _is_remote_image_ref(value: str | Path | None) -> bool:
    """Return whether *value* is an absolute HTTP(S) URL."""
    return is_http_url(value)


def _resolve_image(value: str | Path | None) -> Path | None:
    """Resolve a local path or catalog ID using a deterministic precedence.

    Absolute/user-expanded paths win, then paths relative to the current
    working directory, then bundled catalog IDs. Remote URLs are intentionally
    materialized separately by :func:`_materialize_remote_images`.
    """
    if not value:
        return None
    if _is_remote_image_ref(value):
        return None
    candidate = Path(str(value)).expanduser()
    if candidate.is_file():
        return candidate.resolve()
    cwd_candidate = Path.cwd() / candidate
    if cwd_candidate.is_file():
        return cwd_candidate.resolve()
    return resolve_catalog_image(value)


def _download_remote_image(url: str, directory: Path) -> Path:
    """Delegate safe remote-image materialization to ``helpers.py``."""
    return download_remote_image(url, directory)


def _materialize_remote_images(spec: dict, directory: Path) -> dict:
    """Deep-copy *spec* and replace remote image references with local files."""
    prepared = copy.deepcopy(spec)
    cache: dict[str, Path] = {}
    for slide_spec in prepared.get("slides", []):
        ref = _image_ref(slide_spec)
        if not _is_remote_image_ref(ref):
            continue
        url = str(ref)
        local_path = cache.get(url)
        if local_path is None:
            local_path = _download_remote_image(url, directory)
            cache[url] = local_path
        slide_spec["image"] = str(local_path)
        slide_spec.pop("image_url", None)
    return prepared


def _picture_cover(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    focal_x: float = 0.5,
    focal_y: float = 0.5,
):
    """Add a picture using cover-crop and a normalized focal point.

    ``focal_x`` and ``focal_y`` use the ``0.0..1.0`` image coordinate space.
    Cropping is distributed around that point while preserving aspect ratio.
    """
    with Image.open(path) as im:
        src_w, src_h = im.size
    src_ratio = src_w / max(src_h, 1)
    frame_ratio = w / max(h, 0.001)
    pic = slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    focal_x = max(0.0, min(1.0, focal_x))
    focal_y = max(0.0, min(1.0, focal_y))
    if src_ratio > frame_ratio:
        visible = frame_ratio / src_ratio
        excess = 1.0 - visible
        left = excess * focal_x
        right = excess - left
        pic.crop_left = max(0.0, min(excess, left))
        pic.crop_right = max(0.0, min(excess, right))
    elif src_ratio < frame_ratio:
        visible = src_ratio / frame_ratio
        excess = 1.0 - visible
        top = excess * focal_y
        bottom = excess - top
        pic.crop_top = max(0.0, min(excess, top))
        pic.crop_bottom = max(0.0, min(excess, bottom))
    return pic


def _blank(prs: Presentation, theme: dict, *, dark: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    color = theme["dark"] if dark else theme["background"]
    _rect(slide, 0, 0, W_IN, H_IN, color)
    return slide


def _eyebrow(slide, text: str, theme: dict, *, dark: bool = False):
    if not text:
        return
    color = theme["accent_2"] if dark else theme["accent"]
    _textbox(
        slide,
        text.upper(),
        MARGIN,
        0.34,
        5.5,
        0.28,
        size=10,
        color=color,
        font=theme["font_body"],
        bold=True,
    )


def _header(slide, title: str, theme: dict, *, eyebrow: str = ""):
    _eyebrow(slide, eyebrow, theme)
    y = 0.68 if eyebrow else 0.5
    _textbox(
        slide,
        title,
        MARGIN,
        y,
        W_IN - 2 * MARGIN,
        1.3,
        size=35,
        color=theme["text"],
        font=theme["font_headline"],
        bold=True,
    )


def _footer(
    slide,
    theme: dict,
    page: int,
    *,
    source: str = "",
    dark: bool = False,
):
    text_color = theme["on_dark"] if dark else theme["muted"]
    if source:
        _textbox(
            slide,
            source,
            MARGIN,
            7.12,
            10.8,
            0.18,
            size=9,
            color=text_color,
            font=theme["font_body"],
        )
    _textbox(
        slide,
        f"{page:02d}",
        12.03,
        7.06,
        0.5,
        0.22,
        size=10,
        color=text_color,
        font=theme["font_body"],
        align=PP_ALIGN.RIGHT,
    )


def _slide_title(prs, slide_spec: dict, theme: dict, page: int):
    path = _resolve_image(_image_ref(slide_spec))
    slide = _blank(prs, theme, dark=True)
    if path:
        _picture_cover(
            slide,
            path,
            0,
            0,
            W_IN,
            H_IN,
            focal_x=float(slide_spec.get("focal_x", 0.5)),
            focal_y=float(slide_spec.get("focal_y", 0.5)),
        )
        _rect(slide, 0, 0, W_IN, H_IN, theme["dark"], opacity=0.68)
        _rect(slide, 0, 0, 7.2, H_IN, theme["dark"], opacity=0.28)
    _eyebrow(slide, slide_spec.get("kicker", ""), theme, dark=True)
    _textbox(
        slide,
        slide_spec["title"],
        MARGIN,
        1.72,
        7.7,
        3.15,
        size=50,
        color=theme["on_dark"],
        font=theme["font_headline"],
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if slide_spec.get("subtitle"):
        _textbox(
            slide,
            slide_spec["subtitle"],
            MARGIN,
            5.2,
            6.8,
            0.8,
            size=22,
            color=theme["on_dark"],
            font=theme["font_body"],
        )
    if slide_spec.get("meta"):
        _textbox(
            slide,
            slide_spec["meta"],
            MARGIN,
            6.65,
            7.5,
            0.25,
            size=10,
            color=theme["accent_2"],
            font=theme["font_body"],
            bold=True,
        )
    _footer(slide, theme, page, dark=True)


def _slide_section(prs, slide_spec: dict, theme: dict, page: int):
    slide = _blank(prs, theme, dark=True)
    path = _resolve_image(_image_ref(slide_spec))
    if path:
        _picture_cover(slide, path, 0, 0, W_IN, H_IN)
        _rect(slide, 0, 0, W_IN, H_IN, theme["dark"], opacity=0.72)
    _textbox(
        slide,
        str(slide_spec.get("number", page)).zfill(2),
        MARGIN,
        1.08,
        1.1,
        0.7,
        size=18,
        color=theme["accent_2"],
        font=theme["font_body"],
        bold=True,
    )
    _line(slide, MARGIN, 1.86, 2.0, 1.86, theme["accent"], width=3)
    _textbox(
        slide,
        slide_spec["title"],
        MARGIN,
        2.35,
        9.6,
        1.55,
        size=44,
        color=theme["on_dark"],
        font=theme["font_headline"],
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if slide_spec.get("subtitle"):
        _textbox(
            slide,
            slide_spec["subtitle"],
            MARGIN,
            4.22,
            8.7,
            0.72,
            size=20,
            color=theme["on_dark"],
            font=theme["font_body"],
        )
    _footer(slide, theme, page, dark=True)


def _slide_statement(prs, slide_spec: dict, theme: dict, page: int):
    dark = bool(slide_spec.get("dark", False))
    slide = _blank(prs, theme, dark=dark)
    text = theme["on_dark"] if dark else theme["text"]
    muted = theme["on_dark"] if dark else theme["muted"]
    _eyebrow(slide, slide_spec.get("eyebrow", ""), theme, dark=dark)
    _rect(
        slide,
        MARGIN,
        1.35,
        0.09,
        4.8,
        theme["accent_2"] if dark else theme["accent"],
    )
    _textbox(
        slide,
        slide_spec["title"],
        1.18,
        1.45,
        10.75,
        2.65,
        size=46,
        color=text,
        font=theme["font_headline"],
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if slide_spec.get("body"):
        _textbox(
            slide,
            slide_spec["body"],
            1.2,
            4.4,
            9.25,
            1.05,
            size=21,
            color=muted,
            font=theme["font_body"],
        )
    _footer(
        slide,
        theme,
        page,
        source=slide_spec.get("source", ""),
        dark=dark,
    )


def _slide_bullets(prs, slide_spec: dict, theme: dict, page: int):
    slide = _blank(prs, theme)
    _header(
        slide,
        slide_spec["title"],
        theme,
        eyebrow=slide_spec.get("eyebrow", ""),
    )
    if slide_spec.get("lead"):
        _textbox(
            slide,
            slide_spec["lead"],
            MARGIN,
            2.12,
            10.7,
            0.65,
            size=21,
            color=theme["muted"],
            font=theme["font_body"],
        )
        body_y, body_h = 2.92, 3.62
    else:
        body_y, body_h = 2.18, 4.38
    _add_paragraphs(
        slide,
        slide_spec["bullets"],
        MARGIN,
        body_y,
        10.9,
        body_h,
        size=19,
        color=theme["text"],
        accent=theme["accent"],
        font=theme["font_body"],
        max_items=5,
        spacing=13,
    )
    _footer(slide, theme, page, source=slide_spec.get("source", ""))


def _slide_image(prs, slide_spec: dict, theme: dict, page: int):
    slide = _blank(prs, theme)
    _header(
        slide,
        slide_spec["title"],
        theme,
        eyebrow=slide_spec.get("eyebrow", ""),
    )
    path = _resolve_image(_image_ref(slide_spec))
    if not path:
        raise FileNotFoundError(f"Image not found: {_image_ref(slide_spec)!r}")
    side = slide_spec.get("image_side", "right")
    image_x = 7.18 if side == "right" else 0.0
    text_x = MARGIN if side == "right" else 7.62
    text_w = 5.65 if side == "right" else 4.93
    image_w = W_IN - image_x if side == "right" else 6.66
    _picture_cover(
        slide,
        path,
        image_x,
        2.05,
        image_w,
        5.45,
        focal_x=float(slide_spec.get("focal_x", 0.5)),
        focal_y=float(slide_spec.get("focal_y", 0.5)),
    )
    if slide_spec.get("body"):
        _textbox(
            slide,
            slide_spec["body"],
            text_x,
            2.25,
            text_w,
            1.1,
            size=22,
            color=theme["text"],
            font=theme["font_body"],
            bold=bool(slide_spec.get("body_bold", False)),
        )
        bullets_y = 3.58
    else:
        bullets_y = 2.32
    if slide_spec.get("bullets"):
        _add_paragraphs(
            slide,
            slide_spec["bullets"],
            text_x,
            bullets_y,
            text_w,
            3.0,
            size=17,
            color=theme["text"],
            accent=theme["accent"],
            font=theme["font_body"],
            max_items=4,
            spacing=11,
        )
    _footer(slide, theme, page, source=slide_spec.get("source", ""))


def _slide_metrics(prs, slide_spec: dict, theme: dict, page: int):
    slide = _blank(prs, theme)
    _header(
        slide,
        slide_spec["title"],
        theme,
        eyebrow=slide_spec.get("eyebrow", ""),
    )
    metrics = list(slide_spec["metrics"])[:4]
    n = len(metrics)
    left, right = MARGIN, W_IN - MARGIN
    col_w = (right - left) / max(n, 1)
    for i in range(1, n):
        x = left + col_w * i
        _line(slide, x, 2.32, x, 6.16, theme["line"], width=1.2)
    for i, metric in enumerate(metrics):
        x = left + col_w * i + 0.1
        _textbox(
            slide,
            metric["value"],
            x,
            2.42,
            col_w - 0.22,
            1.2,
            size=48 if n <= 3 else 40,
            color=metric.get("color", theme["accent"]),
            font=theme["font_headline"],
            bold=True,
            valign=MSO_ANCHOR.BOTTOM,
        )
        _textbox(
            slide,
            metric["label"],
            x,
            3.72,
            col_w - 0.22,
            0.72,
            size=18,
            color=theme["text"],
            font=theme["font_body"],
            bold=True,
        )
        if metric.get("detail"):
            _textbox(
                slide,
                metric["detail"],
                x,
                4.62,
                col_w - 0.28,
                1.15,
                size=16,
                color=theme["muted"],
                font=theme["font_body"],
            )
    _footer(slide, theme, page, source=slide_spec.get("source", ""))


def _chart_colors(theme: dict) -> list[RGBColor]:
    return [
        _rgb(theme["accent"]),
        _rgb(theme["accent_2"]),
        _rgb(theme["accent_3"]),
        _rgb(theme["muted"]),
        _rgb(theme["positive"]),
        _rgb(theme["negative"]),
    ]


def _normalize_chart_axis_ids(chart) -> None:
    """Convert signed chart-axis IDs to OOXML's unsigned UInt32 representation.

    Some presentation writers emit negative IDs even though the schema expects
    UInt32. Adding ``2**32`` preserves the bit pattern and improves compatibility
    with strict renderers without changing chart relationships.
    """
    try:
        nodes = chart._chartSpace.xpath(".//c:axId | .//c:crossAx")
    except Exception:
        return
    for node in nodes:
        raw = node.get("val")
        if raw is None:
            continue
        try:
            value = int(raw)
        except ValueError:
            continue
        if value < 0:
            node.set("val", str(value + 2**32))


def _slide_chart(prs, slide_spec: dict, theme: dict, page: int):
    slide = _blank(prs, theme)
    _header(
        slide,
        slide_spec["title"],
        theme,
        eyebrow=slide_spec.get("eyebrow", ""),
    )
    chart_type_name = slide_spec["chart_type"]
    chart_type = CHART_TYPES[chart_type_name]
    categories = [str(c) for c in slide_spec["categories"]]
    series = slide_spec["series"]
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for item in series:
        chart_data.add_series(item["name"], tuple(float(v) for v in item["values"]))

    has_insight = bool(slide_spec.get("insight"))
    chart_w = 8.2 if has_insight else 11.7
    chart_h = 4.15
    chart_shape = slide.shapes.add_chart(
        chart_type,
        Inches(MARGIN),
        Inches(2.32),
        Inches(chart_w),
        Inches(chart_h),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = len(series) > 1 or chart_type_name == "pie"
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = theme["font_body"]
        chart.legend.font.size = Pt(10)

    palette = _chart_colors(theme)
    for i, chart_series in enumerate(chart.series):
        color = palette[i % len(palette)]
        chart_series.format.fill.solid()
        chart_series.format.fill.fore_color.rgb = color
        chart_series.format.line.color.rgb = color
        if chart_type_name == "line":
            chart_series.format.line.width = Pt(2.5)
            chart_series.marker.style = XL_MARKER_STYLE.CIRCLE
            chart_series.marker.size = 7

    if chart_type_name == "pie":
        plot = chart.plots[0]
        plot.has_data_labels = True
        labels = plot.data_labels
        labels.show_percentage = True
        labels.show_category_name = False
        labels.show_legend_key = False
        labels.position = XL_DATA_LABEL_POSITION.BEST_FIT
        labels.font.name = theme["font_body"]
        labels.font.size = Pt(11)
        for point_index, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = palette[point_index % len(palette)]
    else:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = _rgb(theme["line"])
        chart.value_axis.major_gridlines.format.line.width = Pt(0.6)
        chart.value_axis.format.line.fill.background()
        chart.value_axis.tick_labels.font.name = theme["font_body"]
        chart.value_axis.tick_labels.font.size = Pt(9)
        chart.value_axis.tick_labels.font.color.rgb = _rgb(theme["muted"])
        chart.category_axis.tick_labels.font.name = theme["font_body"]
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.category_axis.tick_labels.font.color.rgb = _rgb(theme["muted"])
        chart.category_axis.format.line.color.rgb = _rgb(theme["line"])
        if len(categories) <= 8 and len(series) == 1:
            plot = chart.plots[0]
            plot.has_data_labels = True
            plot.data_labels.position = (
                XL_DATA_LABEL_POSITION.OUTSIDE_END
                if chart_type_name in {"column", "bar"}
                else XL_DATA_LABEL_POSITION.ABOVE
            )
            plot.data_labels.font.name = theme["font_body"]
            plot.data_labels.font.size = Pt(10)
            plot.data_labels.font.color.rgb = _rgb(theme["text"])

    if slide_spec.get("insight"):
        _rect(slide, 9.22, 2.36, 0.07, 3.82, theme["accent"])
        _textbox(
            slide,
            slide_spec.get("insight_label", "КЛЮЧЕВОЙ ВЫВОД").upper(),
            9.58,
            2.46,
            2.72,
            0.35,
            size=10,
            color=theme["accent"],
            font=theme["font_body"],
            bold=True,
        )
        _textbox(
            slide,
            slide_spec["insight"],
            9.58,
            2.94,
            2.72,
            2.8,
            size=18,
            color=theme["text"],
            font=theme["font_body"],
            bold=True,
        )
    if slide_spec.get("unit"):
        _textbox(
            slide,
            slide_spec["unit"],
            MARGIN,
            2.06,
            4.0,
            0.22,
            size=10,
            color=theme["muted"],
            font=theme["font_body"],
        )
    _normalize_chart_axis_ids(chart)
    _footer(slide, theme, page, source=slide_spec.get("source", ""))


def _comparison_column(
    slide,
    block: dict,
    x: float,
    w: float,
    theme: dict,
    *,
    accent: str,
):
    _textbox(
        slide,
        block["heading"],
        x,
        2.3,
        w,
        0.6,
        size=23,
        color=accent,
        font=theme["font_headline"],
        bold=True,
    )
    if block.get("lead"):
        _textbox(
            slide,
            block["lead"],
            x,
            3.12,
            w,
            0.65,
            size=17,
            color=theme["muted"],
            font=theme["font_body"],
        )
        y = 3.92
    else:
        y = 3.08
    _add_paragraphs(
        slide,
        block.get("items", []),
        x,
        y,
        w,
        2.7,
        size=17,
        color=theme["text"],
        accent=accent,
        font=theme["font_body"],
        max_items=4,
        spacing=10,
    )


def _slide_comparison(prs, slide_spec: dict, theme: dict, page: int):
    slide = _blank(prs, theme)
    _header(
        slide,
        slide_spec["title"],
        theme,
        eyebrow=slide_spec.get("eyebrow", ""),
    )
    _line(slide, 6.666, 2.18, 6.666, 6.58, theme["line"], width=1.4)
    _comparison_column(
        slide,
        slide_spec["left"],
        MARGIN,
        5.38,
        theme,
        accent=theme["muted"],
    )
    _comparison_column(
        slide,
        slide_spec["right"],
        7.18,
        5.38,
        theme,
        accent=theme["accent"],
    )
    _footer(slide, theme, page, source=slide_spec.get("source", ""))


def _slide_timeline(prs, slide_spec: dict, theme: dict, page: int, *, process=False):
    slide = _blank(prs, theme)
    _header(
        slide,
        slide_spec["title"],
        theme,
        eyebrow=slide_spec.get("eyebrow", ""),
    )
    steps = list(slide_spec["steps"])[:5]
    n = len(steps)
    x0, x1, y_line = 1.05, 12.15, 3.28
    # Draw connectors before nodes so edges stay behind entities.
    _line(slide, x0, y_line, x1, y_line, theme["line"], width=3)
    positions = [x0 + (x1 - x0) * i / max(n - 1, 1) for i in range(n)]
    if process and n > 1:
        for i in range(n - 1):
            _line(
                slide,
                positions[i] + 0.2,
                y_line,
                positions[i + 1] - 0.2,
                y_line,
                theme["accent"],
                width=2,
            )
    for i, (step, x) in enumerate(zip(steps, positions)):
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x - 0.2),
            Inches(y_line - 0.2),
            Inches(0.4),
            Inches(0.4),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = _rgb(
            theme["accent"] if process else theme["surface"]
        )
        node.line.color.rgb = _rgb(theme["accent"])
        node.line.width = Pt(2)
        label = step.get("label", f"{i + 1:02d}")
        _textbox(
            slide,
            label,
            x - 0.56,
            2.48,
            1.12,
            0.35,
            size=11,
            color=theme["accent"],
            font=theme["font_body"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _textbox(
            slide,
            step["title"],
            x - 0.92,
            3.75,
            1.84,
            0.82,
            size=16,
            color=theme["text"],
            font=theme["font_body"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if step.get("detail"):
            _textbox(
                slide,
                step["detail"],
                x - 0.92,
                4.68,
                1.84,
                1.0,
                size=16,
                color=theme["muted"],
                font=theme["font_body"],
                align=PP_ALIGN.CENTER,
            )
    _footer(slide, theme, page, source=slide_spec.get("source", ""))


def _slide_quote(prs, slide_spec: dict, theme: dict, page: int):
    path = _resolve_image(_image_ref(slide_spec))
    slide = _blank(prs, theme, dark=bool(path))
    dark = bool(path)
    if path:
        _picture_cover(slide, path, 0, 0, W_IN, H_IN)
        _rect(slide, 0, 0, W_IN, H_IN, theme["dark"], opacity=0.72)
    text_color = theme["on_dark"] if dark else theme["text"]
    muted = theme["on_dark"] if dark else theme["muted"]
    _textbox(
        slide,
        "“",
        MARGIN,
        0.72,
        1.2,
        0.68,
        size=62,
        color=theme["accent_2"] if dark else theme["accent"],
        font=theme["font_headline"],
        bold=True,
    )
    _textbox(
        slide,
        slide_spec["quote"],
        1.28,
        1.62,
        10.65,
        3.55,
        size=34,
        color=text_color,
        font=theme["font_headline"],
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if slide_spec.get("attribution"):
        _textbox(
            slide,
            slide_spec["attribution"],
            1.3,
            5.45,
            5.8,
            0.38,
            size=15,
            color=text_color,
            font=theme["font_body"],
            bold=True,
        )
    if slide_spec.get("role"):
        _textbox(
            slide,
            slide_spec["role"],
            1.3,
            5.9,
            6.4,
            0.35,
            size=12,
            color=muted,
            font=theme["font_body"],
        )
    _footer(slide, theme, page, source=slide_spec.get("source", ""), dark=dark)


def _slide_table(prs, slide_spec: dict, theme: dict, page: int):
    slide = _blank(prs, theme)
    _header(
        slide,
        slide_spec["title"],
        theme,
        eyebrow=slide_spec.get("eyebrow", ""),
    )
    columns = [str(c) for c in slide_spec["columns"]]
    rows = [[str(v) for v in row] for row in slide_spec["rows"]]
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        len(columns),
        Inches(MARGIN),
        Inches(2.22),
        Inches(W_IN - 2 * MARGIN),
        Inches(4.48),
    )
    table = table_shape.table
    header_h = 0.58
    table.rows[0].height = Inches(header_h)
    body_h = (4.48 - header_h) / max(len(rows), 1)
    for row_index in range(1, len(table.rows)):
        table.rows[row_index].height = Inches(body_h)
    for col_index, value in enumerate(columns):
        cell = table.cell(0, col_index)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(theme["dark"])
        _set_cell_margins(cell, left=0.12, right=0.12)
        _set_cell_borders(cell, theme["dark_alt"], width_pt=0.45)
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = value
        run.font.name = theme["font_body"]
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = _rgb(theme["on_dark"])
    highlight_row = slide_spec.get("highlight_row")
    for row_index, row_values in enumerate(rows, start=1):
        for col_index, value in enumerate(row_values):
            cell = table.cell(row_index, col_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(
                theme["surface"]
                if highlight_row != row_index - 1
                else theme["background"]
            )
            _set_cell_margins(cell, left=0.12, right=0.12)
            _set_cell_borders(cell, theme["line"], width_pt=0.45)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = (
                PP_ALIGN.LEFT
                if col_index == 0
                else slide_spec.get("numeric_align", PP_ALIGN.RIGHT)
            )
            run = p.add_run()
            run.text = value
            run.font.name = theme["font_body"]
            run.font.size = Pt(12)
            run.font.bold = highlight_row == row_index - 1
            run.font.color.rgb = _rgb(theme["text"])
    _footer(slide, theme, page, source=slide_spec.get("source", ""))


def _slide_summary(prs, slide_spec: dict, theme: dict, page: int):
    slide = _blank(prs, theme)
    _header(
        slide,
        slide_spec["title"],
        theme,
        eyebrow=slide_spec.get("eyebrow", "ИТОГ"),
    )
    _textbox(
        slide,
        slide_spec["takeaway"],
        MARGIN,
        2.1,
        11.4,
        1.25,
        size=31,
        color=theme["text"],
        font=theme["font_headline"],
        bold=True,
    )
    _line(slide, MARGIN, 3.48, 12.55, 3.48, theme["line"], width=1.2)
    _add_paragraphs(
        slide,
        slide_spec["actions"],
        MARGIN,
        3.82,
        11.1,
        2.55,
        size=19,
        color=theme["text"],
        accent=theme["accent"],
        font=theme["font_body"],
        numbered=True,
        max_items=4,
        spacing=13,
    )
    _footer(slide, theme, page, source=slide_spec.get("source", ""))


def _slide_closing(prs, slide_spec: dict, theme: dict, page: int):
    slide = _blank(prs, theme, dark=True)
    path = _resolve_image(_image_ref(slide_spec))
    if path:
        _picture_cover(
            slide,
            path,
            0,
            0,
            W_IN,
            H_IN,
            focal_x=float(slide_spec.get("focal_x", 0.5)),
            focal_y=float(slide_spec.get("focal_y", 0.5)),
        )
        _rect(slide, 0, 0, W_IN, H_IN, theme["dark"], opacity=0.67)
    _line(slide, MARGIN, 1.25, 2.25, 1.25, theme["accent_2"], width=4)
    _textbox(
        slide,
        slide_spec["title"],
        MARGIN,
        2.12,
        8.1,
        2.2,
        size=45,
        color=theme["on_dark"],
        font=theme["font_headline"],
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if slide_spec.get("subtitle"):
        _textbox(
            slide,
            slide_spec["subtitle"],
            MARGIN,
            4.68,
            7.5,
            0.9,
            size=20,
            color=theme["on_dark"],
            font=theme["font_body"],
        )
    if slide_spec.get("contact"):
        _textbox(
            slide,
            slide_spec["contact"],
            MARGIN,
            6.35,
            7.7,
            0.3,
            size=11,
            color=theme["accent_2"],
            font=theme["font_body"],
            bold=True,
        )
    _footer(slide, theme, page, dark=True)


BUILDERS = {
    "title": _slide_title,
    "section": _slide_section,
    "statement": _slide_statement,
    "bullets": _slide_bullets,
    "image": _slide_image,
    "metrics": _slide_metrics,
    "chart": _slide_chart,
    "comparison": _slide_comparison,
    "timeline": lambda prs, spec, theme, page: _slide_timeline(
        prs, spec, theme, page, process=False
    ),
    "process": lambda prs, spec, theme, page: _slide_timeline(
        prs, spec, theme, page, process=True
    ),
    "quote": _slide_quote,
    "table": _slide_table,
    "summary": _slide_summary,
    "closing": _slide_closing,
}


def _missing(value: Any) -> bool:
    if value is None:
        return True
    if isinstance(value, str):
        return not value.strip()
    if isinstance(value, (list, dict, tuple)):
        return len(value) == 0
    return False


def lint_spec(spec: dict) -> dict[str, list[str]]:
    """Validate a presentation spec without mutating it or writing files.

    Args:
        spec: Dictionary described in ``SPEC_REFERENCE.md``.

    Returns:
        A dictionary with ``errors`` and ``warnings`` lists. Errors block
        rendering; warnings flag density, sourcing, or narrative concerns.
        Remote image URLs are syntax-checked here and downloaded only by
        :func:`build_from_spec`.
    """
    errors: list[str] = []
    warnings: list[str] = []
    # Phase 1: validate deck-level structure and resolve the chosen theme.
    if not isinstance(spec, dict):
        return {"errors": ["spec must be a dict"], "warnings": []}
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        return {"errors": ["spec['slides'] must be a non-empty list"], "warnings": []}
    if len(slides) > 40:
        warnings.append("The deck has more than 40 slides; consider splitting it.")
    if isinstance(spec.get("theme"), str):
        try:
            _load_theme(spec.get("theme"))
        except Exception as exc:
            errors.append(str(exc))
    previous_type = None
    repeated = 1
    # Phase 2: validate each slide's required fields and type-specific invariants.
    for index, slide in enumerate(slides, start=1):
        prefix = f"Slide {index}"
        if not isinstance(slide, dict):
            errors.append(f"{prefix}: slide must be an object")
            continue
        slide_type = slide.get("type")
        if slide_type not in SLIDE_TYPES:
            errors.append(
                f"{prefix}: unknown type {slide_type!r}; available: {sorted(SLIDE_TYPES)}"
            )
            continue
        for field in REQUIRED_FIELDS[slide_type]:
            if _missing(slide.get(field)):
                errors.append(f"{prefix} ({slide_type}): missing non-empty {field!r}")
        title = slide.get("title")
        if isinstance(title, str) and len(title) > 82:
            warnings.append(f"{prefix}: title is {len(title)} characters; shorten it.")
        if slide_type == "bullets":
            bullets = slide.get("bullets", [])
            if len(bullets) > 5:
                warnings.append(f"{prefix}: only the first 5 bullets will be rendered.")
            for item in bullets:
                if len(str(item)) > 170:
                    warnings.append(f"{prefix}: a bullet is unusually long.")
        if slide_type in {"timeline", "process"} and len(slide.get("steps", [])) > 5:
            warnings.append(f"{prefix}: only the first 5 steps will be rendered.")
        if slide_type == "metrics":
            metrics = slide.get("metrics", [])
            if not 2 <= len(metrics) <= 4:
                errors.append(f"{prefix}: metrics must contain 2–4 items.")
            for metric in metrics:
                if (
                    not isinstance(metric, dict)
                    or _missing(metric.get("value"))
                    or _missing(metric.get("label"))
                ):
                    errors.append(f"{prefix}: every metric needs value and label.")
        if slide_type not in IMAGE_SLIDE_TYPES and (
            not _missing(slide.get("image")) or not _missing(slide.get("image_url"))
        ):
            errors.append(
                f"{prefix} ({slide_type}): this slide type has no image field."
            )
        if slide_type in IMAGE_SLIDE_TYPES:
            has_image = not _missing(slide.get("image"))
            has_image_url = not _missing(slide.get("image_url"))
            if has_image and has_image_url:
                errors.append(
                    f"{prefix}: use only one of 'image' or 'image_url', not both."
                )
            image_ref = _image_ref(slide)
            if slide_type == "image" and _missing(image_ref):
                errors.append(
                    f"{prefix} (image): provide a catalog ID/local path in "
                    f"'image' or an HTTP(S) URL in 'image_url'."
                )
            if has_image_url and not _is_remote_image_ref(slide.get("image_url")):
                errors.append(f"{prefix}: image_url must be an absolute HTTP(S) URL.")
            if image_ref and not _is_remote_image_ref(image_ref):
                if not _resolve_image(image_ref):
                    errors.append(f"{prefix}: image {image_ref!r} was not found.")
        if slide_type == "image":
            if slide.get("image_side", "right") not in {"left", "right"}:
                errors.append(f"{prefix}: image_side must be left or right.")
        if slide_type == "chart":
            chart_type = slide.get("chart_type")
            if chart_type not in CHART_TYPES:
                errors.append(
                    f"{prefix}: chart_type must be one of {sorted(CHART_TYPES)}."
                )
            categories = slide.get("categories", [])
            series = slide.get("series", [])
            if chart_type == "pie" and len(series) != 1:
                errors.append(f"{prefix}: pie charts require exactly one series.")
            if chart_type == "pie" and len(categories) > 6:
                warnings.append(f"{prefix}: pie chart has more than 6 segments.")
            if len(categories) > 14:
                warnings.append(f"{prefix}: more than 14 categories may be unreadable.")
            for series_index, item in enumerate(series, start=1):
                if not isinstance(item, dict) or _missing(item.get("name")):
                    errors.append(f"{prefix}: series {series_index} needs a name.")
                    continue
                values = item.get("values")
                if not isinstance(values, list) or len(values) != len(categories):
                    errors.append(
                        f"{prefix}: series {series_index} values must match categories."
                    )
                    continue
                for value in values:
                    try:
                        float(value)
                    except (TypeError, ValueError):
                        errors.append(
                            f"{prefix}: series {series_index} contains non-numeric data."
                        )
                        break
            if not slide.get("source"):
                warnings.append(f"{prefix}: add a source for chart data.")
        if slide_type == "table":
            columns = slide.get("columns", [])
            rows = slide.get("rows", [])
            if len(columns) > 8 or len(rows) > 9:
                warnings.append(f"{prefix}: table may be too dense.")
            for row_index, row in enumerate(rows, start=1):
                if len(row) != len(columns):
                    errors.append(
                        f"{prefix}: row {row_index} has {len(row)} cells, "
                        f"expected {len(columns)}."
                    )
        if slide_type == previous_type:
            repeated += 1
            if repeated == 4:
                warnings.append(
                    f"{prefix}: four consecutive {slide_type!r} slides reduce visual rhythm."
                )
        else:
            previous_type = slide_type
            repeated = 1
    # Phase 3: apply deck-level narrative checks after all slides are known.
    if slides and slides[0].get("type") != "title":
        warnings.append("The first slide is not a title slide.")
    if slides and slides[-1].get("type") != "closing":
        warnings.append("The last slide is not a closing slide.")
    return {"errors": errors, "warnings": warnings}


def validate_spec(spec: dict) -> list[str]:
    """Return only blocking errors from :func:`lint_spec`.

    Args:
        spec: Presentation dictionary to validate without mutation.

    Returns:
        List of blocking error messages; an empty list means the spec may be
        passed to :func:`build_from_spec`.

    Use this compatibility wrapper when warnings are intentionally ignored.
    New callers should prefer :func:`lint_spec` so they can surface both lists.
    """
    return lint_spec(spec)["errors"]


def _safe_filename(value: str) -> str:
    return sanitize_filename(value)


def build_from_spec(spec: dict, out_dir: str | Path | None = None) -> str:
    """Build a presentation and return the absolute output path.

    Args:
        spec: Presentation dictionary described in
            ``SPEC_REFERENCE.md``. The input is not mutated.
        out_dir: Destination directory. Defaults to the current directory and
            is created when missing.

    Returns:
        Absolute path to the written ``.pptx`` file as a string.

    Raises:
        ValueError: If :func:`lint_spec` reports blocking errors or content
            cannot fit its assigned text geometry at the designed font size.
        RuntimeError: If a remote image cannot be materialized or the output is
            missing/empty after saving, or structural QA finds overflow,
            out-of-bounds objects, or unintended overlaps.
        OSError: If the destination cannot be created or written.

    Notes:
        ``filename`` is sanitized to a portable basename. An existing file with
        the same name is replaced. Remote images are temporary and embedded
        before their downloaded copies are deleted.
    """
    assert_bundle_complete()
    report = lint_spec(spec)
    if report["errors"]:
        raise ValueError(
            "Invalid presentation spec:\n- " + "\n- ".join(report["errors"])
        )
    destination = Path(out_dir or Path.cwd()).resolve()
    destination.mkdir(parents=True, exist_ok=True)
    filename = _safe_filename(spec.get("filename", "presentation.pptx"))
    path = destination / filename
    with tempfile.TemporaryDirectory(prefix="presentation-images-") as tmp:
        prepared = _materialize_remote_images(spec, Path(tmp))
        theme = _load_theme(prepared.get("theme"))
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        for page, slide_spec in enumerate(prepared["slides"], start=1):
            BUILDERS[slide_spec["type"]](prs, slide_spec, theme, page)
        prs.save(path)
    if not path.exists() or path.stat().st_size == 0:
        raise RuntimeError(f"Presentation was not written correctly: {path}")
    try:
        from .validate_deck import inspect_deck
    except ImportError:
        from validate_deck import inspect_deck

    qa = inspect_deck(path)
    if qa["issues"]:
        raise RuntimeError(
            f"Presentation failed structural QA and must not be delivered: {path}\n- "
            + "\n- ".join(qa["issues"])
        )
    return str(path)


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Validate or build a presentation from a JSON spec.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Examples:\n"
            "  python3 builder.py spec.json --lint-only\n"
            "  python3 builder.py spec.json --output-dir ./output\n\n"
            "Exit status: 0 on success; 2 when the spec is invalid."
        ),
    )
    parser.add_argument("spec", type=Path, help="Path to the presentation JSON spec.")
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=Path.cwd(),
        help="Directory for the PPTX (default: current directory).",
    )
    parser.add_argument(
        "--lint-only",
        action="store_true",
        help="Validate the spec without writing a PPTX.",
    )
    args = parser.parse_args()

    payload = _load_json(args.spec)
    result = lint_spec(payload)
    for warning in result["warnings"]:
        print(f"WARNING: {warning}")
    if result["errors"]:
        for error in result["errors"]:
            print(f"ERROR: {error}")
        raise SystemExit(2)
    if args.lint_only:
        print("Spec is valid.")
    else:
        print(build_from_spec(payload, args.output_dir))
