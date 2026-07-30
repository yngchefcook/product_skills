#!/usr/bin/env python3
"""Run deterministic structural checks against a generated PPTX.

The validator detects out-of-bounds shapes, likely text overflow, suspicious
text/chart overlaps, empty slides, and very small text. It is intentionally
heuristic: a clean report does not replace rendering and visual inspection.

Public API:
    ``inspect_deck(path)`` returns a structured report without printing.

CLI example:
    python3 validate_deck.py presentation.pptx --strict
"""

from __future__ import annotations

import argparse
import math
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# PowerPoint stores geometry in English Metric Units (EMU). Tolerances absorb
# tiny rounding differences introduced by OOXML writers and renderers.
EMU_PER_INCH = 914400
EDGE_TOLERANCE = 6000
OVERLAP_TOLERANCE = 18000

__all__ = ["inspect_deck"]


def _bbox(shape):
    try:
        left = int(shape.left)
        top = int(shape.top)
        width = int(shape.width)
        height = int(shape.height)
    except (AttributeError, TypeError, ValueError):
        return None
    if width <= 0 or height <= 0:
        return None
    return left, top, left + width, top + height


def _overlap(a, b):
    width = min(a[2], b[2]) - max(a[0], b[0])
    height = min(a[3], b[3]) - max(a[1], b[1])
    return width > OVERLAP_TOLERANCE and height > OVERLAP_TOLERANCE


def _text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return " ".join(shape.text_frame.text.split())


def _raw_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text_frame.text


def _text_font_sizes(shape) -> list[float]:
    sizes: list[float] = []
    if not getattr(shape, "has_text_frame", False):
        return sizes
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
    return sizes


def _estimate_text_capacity(shape) -> tuple[int, int]:
    """Return estimated required and available line counts.

    The constants approximate average glyph width and line height. The result
    is a conservative warning heuristic, not a typography engine.
    """
    text = _raw_text(shape)
    if not text:
        return 0, 1
    bbox = _bbox(shape)
    if not bbox:
        return 0, 1
    width_in = (bbox[2] - bbox[0]) / EMU_PER_INCH
    height_in = (bbox[3] - bbox[1]) / EMU_PER_INCH
    sizes = _text_font_sizes(shape)
    font_pt = max(sizes) if sizes else 18
    avg_char_in = max(0.055, font_pt / 72 * 0.52)
    chars_per_line = max(4, int(width_in / avg_char_in))
    required = 0
    for part in text.splitlines() or [""]:
        words = part.split()
        if not words:
            required += 1
            continue
        current = 0
        for word in words:
            word_length = len(word)
            if current == 0:
                required += max(1, math.ceil(word_length / chars_per_line))
                current = word_length % chars_per_line
                continue
            needed = 1 + word_length
            if current + needed <= chars_per_line:
                current += needed
            else:
                required += max(1, math.ceil(word_length / chars_per_line))
                current = word_length % chars_per_line
    line_height_in = font_pt / 72 * 1.16
    available = max(1, int(height_in / max(line_height_in, 0.01)))
    return required, available


def inspect_deck(path: str | Path) -> dict:
    """Inspect a PPTX and return structural issues and non-blocking warnings.

    Args:
        path: Path to an existing ``.pptx`` file.

    Returns:
        Dictionary with:
        ``path`` (absolute path), ``slides`` (count), ``issues`` (blocking
        geometry/overflow findings), and ``warnings`` (advisory findings).

    Notes:
        Overlap detection is pairwise and may produce false positives for
        intentionally layered designs. Render every slide before delivery.
    """
    deck_path = Path(path)
    prs = Presentation(deck_path)
    issues: list[str] = []
    warnings: list[str] = []
    slide_w, slide_h = int(prs.slide_width), int(prs.slide_height)
    for slide_index, slide in enumerate(prs.slides, start=1):
        items = []
        nonempty_text = 0
        for shape_index, shape in enumerate(slide.shapes, start=1):
            bbox = _bbox(shape)
            if bbox:
                if (
                    bbox[0] < -EDGE_TOLERANCE
                    or bbox[1] < -EDGE_TOLERANCE
                    or bbox[2] > slide_w + EDGE_TOLERANCE
                    or bbox[3] > slide_h + EDGE_TOLERANCE
                ):
                    issues.append(
                        f"Slide {slide_index}, shape {shape_index}: outside slide bounds."
                    )
            text = _text(shape)
            if text:
                nonempty_text += 1
                required, available = _estimate_text_capacity(shape)
                if required > available:
                    issues.append(
                        f"Slide {slide_index}: likely text overflow in "
                        f"“{text[:55]}” ({required} estimated lines, {available} fit)."
                    )
                sizes = _text_font_sizes(shape)
                if sizes and min(sizes) < 9:
                    warnings.append(
                        f"Slide {slide_index}: text below 9 pt in “{text[:55]}”."
                    )
            if bbox:
                items.append((shape, bbox, text))
        if nonempty_text == 0:
            warnings.append(f"Slide {slide_index}: no visible text.")

        # Pairwise checks are acceptable here because slide shape counts are
        # small; the bottom metadata strip is intentionally allowed to overlap.
        for i in range(len(items)):
            shape_a, box_a, text_a = items[i]
            if not text_a:
                continue
            for j in range(i + 1, len(items)):
                shape_b, box_b, text_b = items[j]
                if not text_b:
                    continue
                # Source/footer text and slide numbers share the bottom strip by design.
                if box_a[1] > slide_h * 0.92 and box_b[1] > slide_h * 0.92:
                    continue
                if _overlap(box_a, box_b):
                    issues.append(
                        f"Slide {slide_index}: overlapping text boxes "
                        f"“{text_a[:38]}” / “{text_b[:38]}”."
                    )

        for shape, bbox, text in items:
            if int(shape.shape_type) == int(MSO_SHAPE_TYPE.CHART):
                for other, other_bbox, other_text in items:
                    if other is shape or not other_text:
                        continue
                    if _overlap(bbox, other_bbox):
                        issues.append(
                            f"Slide {slide_index}: chart overlaps text "
                            f"“{other_text[:45]}”."
                        )
    return {
        "path": str(deck_path.resolve()),
        "slides": len(prs.slides),
        "issues": issues,
        "warnings": warnings,
    }


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Run structural quality checks against a PPTX deck.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Example:\n"
            "  python3 validate_deck.py presentation.pptx --strict\n\n"
            "Exit status: 0 when usable; 2 for detected issues in strict mode."
        ),
    )
    parser.add_argument("deck", type=Path, help="Path to the PPTX file to inspect.")
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Exit with status 2 when any blocking issue is detected.",
    )
    args = parser.parse_args()
    report = inspect_deck(args.deck)
    print(f"Deck: {report['path']}")
    print(f"Slides: {report['slides']}")
    for warning in report["warnings"]:
        print(f"WARNING: {warning}")
    for issue in report["issues"]:
        print(f"ERROR: {issue}")
    if report["issues"]:
        print(f"Result: {len(report['issues'])} issue(s)")
        if args.strict:
            raise SystemExit(2)
    else:
        print("Result: OK")


if __name__ == "__main__":
    main()
